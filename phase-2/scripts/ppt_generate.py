import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

# ========= PATH =========
BASE_DIR = Path(__file__).resolve().parent
ROOT_DIR = BASE_DIR.parent / "projects"
TEMPLATE = BASE_DIR.parent / "template.pptx"
URL_FILE = BASE_DIR.parent / "urls.txt"


# ---------- Extract project name ----------
def extract_project_name(prompt_text: str, fallback: str):
    match = re.search(r"\*\*Name:\*\*\s*(.+)", prompt_text)
    if match:
        name = match.group(1).strip()
        return re.sub(r'[\\/*?:"<>|]', "", name)
    return fallback


# ---------- Parse URLs ----------
def parse_urls(file_path: Path):
    mapping = {}

    if not file_path.exists():
        return mapping

    lines = file_path.read_text(encoding="utf-8").splitlines()

    slug = None

    for line in lines:
        line = line.strip()
        if not line:
            continue

        if line.lower().startswith("slug"):
            parts = line.split(":", 1)
            if len(parts) > 1:
                slug = parts[1].strip()

        elif line.lower().startswith("url") and slug:
            if ":" in line:
                url = line.split(":", 1)[1].strip()
            else:
                parts = line.split()
                if len(parts) < 2:
                    continue
                url = parts[1].strip()

            if not url.startswith("http"):
                url = f"https://{url}"

            mapping[slug] = url
            slug = None

    return mapping


# ---------- Parse ppt.md ----------
def parse_ppt_md(text: str):
    sections = {
        "Brief": "",
        "Opportunities": "",
        "Features": "",
        "Technologies": ""
    }

    current = None

    for line in text.splitlines():
        line = line.strip()

        if line.startswith("# Brief"):
            current = "Brief"
        elif line.startswith("# Opportunities"):
            current = "Opportunities"
        elif line.startswith("# Features"):
            current = "Features"
        elif line.startswith("# Technologies"):
            current = "Technologies"
        elif current:
            sections[current] += line + "\n"

    return sections


# ---------- Clean markdown ----------
def clean_line(line: str):
    line = line.strip()

    line = re.sub(r"^#+\s*", "", line)
    line = line.replace("**", "")

    if line.startswith("- "):
        line = line[2:]

    return line.strip()


# ---------- Clean technologies ----------
def clean_technologies(text: str):
    cleaned = []

    for line in text.splitlines():
        raw = line.strip()

        if not raw:
            continue

        # STOP when unwanted section starts
        if any(raw.lower().startswith(x) for x in [
            "no markdown",
            "no code",
            "clean bullet",
            "keep concise"
        ]):
            break

        if raw.startswith("--- FILE:"):
            break

        line = clean_line(raw)

        if not line:
            continue

        cleaned.append(line)

    return cleaned


# ---------- Write section ----------
def set_section(slide, heading, content_lines, font_size=18):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        if heading.lower() in shape.text.lower():
            tf = shape.text_frame
            tf.clear()

            p = tf.paragraphs[0]
            p.text = f"{heading}:"
            p.font.bold = True
            p.font.size = Pt(font_size + 4)

            for line in content_lines:
                p = tf.add_paragraph()
                p.text = f"• {line}"
                p.font.size = Pt(font_size)

            break


# ---------- Generate PPT ----------
def generate_ppt(project_path: Path):
    ppt_md = project_path / "ppt.md"
    prompt_md = project_path / "prompt.md"

    if not ppt_md.exists():
        print(f"Skipped: {project_path.name}")
        return

    prs = Presentation(TEMPLATE)

    content = parse_ppt_md(ppt_md.read_text(encoding="utf-8"))

    # Clean sections
    brief = [clean_line(x) for x in content["Brief"].splitlines() if clean_line(x)]
    opp = [clean_line(x) for x in content["Opportunities"].splitlines() if clean_line(x)]
    feat = [clean_line(x) for x in content["Features"].splitlines() if clean_line(x)]
    tech = clean_technologies(content["Technologies"])

    # ✅ Extract project name from prompt.md
    project_name = project_path.name
    if prompt_md.exists():
        prompt_text = prompt_md.read_text(encoding="utf-8")
        project_name = extract_project_name(prompt_text, project_name)

    # Slides
    set_section(prs.slides[1], "Brief", brief, 20)
    set_section(prs.slides[2], "Opportunities", opp, 16)
    set_section(prs.slides[3], "Features", feat, 16)
    set_section(prs.slides[4], "Technologies", tech, 20)

    # ✅ Save using project name (NOT folder name)
    output = project_path / f"{project_name}.pptx"
    prs.save(output)

    print(f"Generated: {output}")


# ---------- Main ----------
def main():
    projects = [p for p in ROOT_DIR.iterdir() if p.is_dir()]

    for p in projects:
        generate_ppt(p)

    print("\nDone.")


if __name__ == "__main__":
    main()