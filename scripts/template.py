import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

ROOT_DIR = Path("projects")
TEMPLATE = Path("template.pptx")


# ---------- Extract project name ----------
def extract_project_name(prompt_text: str, fallback: str):
    match = re.search(r"\*\*Name:\*\*\s*(.+)", prompt_text)
    if match:
        name = match.group(1).strip()
        return re.sub(r'[\\/*?:"<>|]', "", name)
    return fallback


# ---------- Parse ppt.md ----------
def parse_ppt_md(text: str):
    sections = {
        "Brief": "",
        "Opportunities": "",
        "Features": "",
        "Technologies": ""
    }

    current = None

    for line in text.splitlines():
        line = line.strip()

        if line.startswith("# Brief"):
            current = "Brief"
        elif line.startswith("# Opportunities"):
            current = "Opportunities"
        elif line.startswith("# Features"):
            current = "Features"
        elif line.startswith("# Technologies"):
            current = "Technologies"
        elif current:
            sections[current] += line + "\n"

    return sections


# ---------- Clean section (REMOVE CODE LEAK) ----------
def clean_section_content(text: str):
    lines = text.splitlines()
    cleaned = []

    for line in lines:
        line = line.strip()

        # STOP if code / file dump starts
        if line.startswith("--- FILE:"):
            break
        if line.startswith("```"):
            break

        cleaned.append(line)

    return "\n".join(cleaned)


# ---------- Clean markdown ----------
def clean_markdown_line(line: str):
    line = line.strip()

    # remove markdown headings (##, ###, etc.)
    if line.startswith("#"):
        line = re.sub(r"^#+\s*", "", line)

    # skip code-like lines
    if (
        line.startswith("import ") or
        line.startswith("export ") or
        line.startswith("const ") or
        line.startswith("function ") or
        "{" in line and "}" in line
    ):
        return ""

    # remove bullet prefix
    if line.startswith("- "):
        line = line[2:]

    # remove bold markdown
    line = line.replace("**", "")

    return line.strip()


# ---------- Write formatted text ----------
def set_section_text(slide, heading, content, is_bullets=False, font_size=18):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        tf = shape.text_frame

        if heading.lower() in tf.text.lower():

            tf.clear()

            # ---- Heading ----
            p = tf.paragraphs[0]
            p.text = f"{heading}:"
            p.font.bold = True
            p.font.size = Pt(font_size + 4)

            # ---- Content ----
            lines = content.split("\n")

            for line in lines:
                line = clean_markdown_line(line)
                if not line:
                    continue

                p = tf.add_paragraph()

                if is_bullets:
                    p.text = f"• {line}"
                else:
                    p.text = line

                p.level = 0
                p.font.size = Pt(font_size)

            break


# ---------- Generate PPT ----------
def generate_ppt(project_path: Path):
    ppt_md = project_path / "ppt.md"
    overview_md = project_path / "overview.md"
    prompt_md = project_path / "prompt.md"

    if not ppt_md.exists():
        print(f"Skipped (no ppt.md): {project_path.name}")
        return

    prs = Presentation(TEMPLATE)

    content = parse_ppt_md(ppt_md.read_text(encoding="utf-8"))

    overview = (
        overview_md.read_text(encoding="utf-8")
        if overview_md.exists()
        else content["Brief"]
    )

    # Extract project name
    project_name = project_path.name
    if prompt_md.exists():
        prompt_text = prompt_md.read_text(encoding="utf-8")
        project_name = extract_project_name(prompt_text, project_name)

    try:
        # Slide 2 → Brief
        # Brief → larger text
        set_section_text(
            prs.slides[1],
            "Brief",
            overview,
            is_bullets=False,
            font_size=21
        )

        # Opportunities
        set_section_text(
            prs.slides[2],
            "Opportunities",
            content["Opportunities"],
            is_bullets=True,
            font_size=16
        )

        # Features
        set_section_text(
            prs.slides[3],
            "Features",
            clean_section_content(content["Features"]),
            is_bullets=True,
            font_size=16
        )

        # Technologies → larger text
        set_section_text(
            prs.slides[4],
            "Technologies",
            clean_section_content(content["Technologies"]),
            is_bullets=True,
            font_size=21
        )

    except IndexError:
        print(f"Template mismatch: {project_name}")
        return

    output_path = project_path / f"{project_name}.pptx"
    prs.save(output_path)

    print(f"Generated: {output_path}")


# ---------- Main ----------
def main():
    if not TEMPLATE.exists():
        raise FileNotFoundError("template.pptx not found in root")

    projects = [p for p in ROOT_DIR.iterdir() if p.is_dir()]

    print(f"Processing {len(projects)} projects...\n")

    for project in projects:
        generate_ppt(project)

    print("\nAll PPTs generated successfully.")


if __name__ == "__main__":
    main()